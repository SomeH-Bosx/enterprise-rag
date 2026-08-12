"""PPTX / legacy PPT loaders.

Legacy .ppt is converted to .pptx first (explicit Office/LibreOffice step),
then parsed with python-pptx.

Step3.6: table shapes → Markdown. Empty slides are skipped with a log
(OCR render for PPTX is not wired; PDF remains the OCR target).
"""

from __future__ import annotations

from pathlib import Path

from langchain_core.documents import Document

from src.config.logging import get_logger
from src.config.settings import Settings, get_settings
from src.ingestion.loaders._common import docs_or_fail, ensure_file
from src.ingestion.office_convert import prepare_for_load
from src.ingestion.tables import pptx_table_to_markdown
from src.services.exceptions import IngestError

logger = get_logger("ppt_loader")


def load_pptx(path: str | Path, settings: Settings | None = None) -> list[Document]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover
        raise IngestError("python-pptx is required for .pptx support") from exc

    cfg = settings or get_settings()
    pptx_path = ensure_file(path, expected_suffixes=(".pptx",))
    presentation = Presentation(str(pptx_path))
    docs: list[Document] = []

    for idx, slide in enumerate(presentation.slides):
        parts: list[str] = []
        table_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                if cfg.enable_table_serialization:
                    md = pptx_table_to_markdown(shape.table)
                    if md:
                        parts.append(md)
                        table_count += 1
                else:
                    for row in shape.table.rows:
                        cells = [
                            (cell.text_frame.text if cell.text_frame is not None else "").strip()
                            for cell in row.cells
                        ]
                        line = " | ".join(c for c in cells if c)
                        if line:
                            parts.append(line)
                    table_count += 1
                continue

            # Skip pure table shapes already handled; avoid double-counting text
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue

            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    parts.append(text)

        body = "\n\n".join(parts).strip()
        if not body:
            logger.info(
                "pptx_slide_empty_skipped",
                slide=idx + 1,
                reason="no_extractable_text",
                ocr_status="skipped_not_needed",
                note="PPTX OCR render not enabled; PDF pages are OCR targets",
            )
            continue

        docs.append(
            Document(
                page_content=body,
                metadata={
                    "source": str(pptx_path),
                    "filename": pptx_path.name,
                    "file_type": "pptx",
                    "page": idx + 1,
                    "slide": idx + 1,
                    "table_count": table_count,
                    "ocr_status": "skipped_not_needed",
                },
            )
        )
    return docs_or_fail(docs, label="PPTX")


def load_ppt(path: str | Path, settings: Settings | None = None) -> list[Document]:
    """
    Legacy .ppt: explicit convert → .pptx, then load_pptx.
    Prefer Microsoft PowerPoint COM; fallback LibreOffice.
    """
    ppt_path = ensure_file(path, expected_suffixes=(".ppt",))
    outcome = prepare_for_load(ppt_path)
    docs = load_pptx(outcome.load_path, settings=settings)
    for d in docs:
        meta = dict(d.metadata or {})
        meta["filename"] = ppt_path.name
        meta["source"] = str(ppt_path)
        meta["file_type"] = "ppt"
        meta["loader"] = f"convert:{outcome.engine}->pptx"
        meta["converted"] = outcome.converted
        meta["convert_engine"] = outcome.engine
        d.metadata = meta
    return docs_or_fail(docs, label="PPT")
