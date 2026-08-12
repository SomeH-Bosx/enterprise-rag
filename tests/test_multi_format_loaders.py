"""Step2 multi-format document loader tests."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from src.ingestion.loaders import (
    SUPPORTED_SUFFIXES,
    detect_file_type,
    is_supported,
    load_document,
    load_docx,
    load_md,
    load_pptx,
    load_txt,
)
from src.services.exceptions import IngestError


def test_supported_suffixes_include_required_formats():
    required = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".md", ".txt"}
    assert required.issubset(set(SUPPORTED_SUFFIXES))


def test_detect_and_is_supported(tmp_path: Path):
    p = tmp_path / "a.md"
    p.write_text("# hi", encoding="utf-8")
    assert detect_file_type(p) == "md"
    assert is_supported(p)
    assert not is_supported(tmp_path / "a.xlsx")


def test_load_txt_and_md(tmp_path: Path):
    txt = tmp_path / "note.txt"
    txt.write_text("年假政策：每年 15 天。", encoding="utf-8")
    md = tmp_path / "note.md"
    md.write_text("# Policy\n\n年假 15 天。", encoding="utf-8")

    txt_docs = load_txt(txt)
    md_docs = load_md(md)
    assert "15" in txt_docs[0].page_content
    assert "Policy" in md_docs[0].page_content
    assert txt_docs[0].metadata["file_type"] == "txt"
    assert md_docs[0].metadata["file_type"] == "md"


def test_load_docx(tmp_path: Path):
    path = tmp_path / "handbook.docx"
    doc = DocxDocument()
    doc.add_heading("Employee Handbook", level=1)
    doc.add_paragraph("Annual leave is 15 days per year.")
    doc.save(path)

    docs = load_docx(path)
    assert len(docs) == 1
    assert "15 days" in docs[0].page_content
    assert docs[0].metadata["file_type"] == "docx"


def test_load_pptx(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Product SLO"
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = box.text_frame
    tf.text = "p95 latency under 200 milliseconds"
    prs.save(path)

    docs = load_pptx(path)
    blob = "\n".join(d.page_content for d in docs)
    assert "200" in blob or "SLO" in blob or "Product" in blob
    assert docs[0].metadata["file_type"] == "pptx"


def test_load_document_dispatcher(tmp_path: Path):
    path = tmp_path / "readme.txt"
    path.write_text("dispatcher works", encoding="utf-8")
    docs = load_document(path)
    assert docs[0].page_content.startswith("dispatcher")
    assert docs[0].metadata["filename"] == "readme.txt"


def test_unsupported_type_raises(tmp_path: Path):
    path = tmp_path / "data.xlsx"
    path.write_bytes(b"fake")
    with pytest.raises(IngestError, match="Unsupported"):
        load_document(path)
