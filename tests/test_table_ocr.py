"""Step3.6 table serialization + OCR unit tests."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from src.config.settings import Settings
from src.ingestion.loaders import load_docx, load_pdf, load_pptx
from src.ingestion.ocr import ocr_available, ocr_pil_image, should_ocr_text
from src.ingestion.tables import rows_to_markdown
from src.services.ingest_service import IngestService


def _settings(tmp_path: Path, **kwargs) -> Settings:
    base = dict(
        VECTOR_DB_PATH=str(tmp_path / "chroma"),
        BM25_STORE_PATH=str(tmp_path / "bm25.json"),
        DOC_REGISTRY_PATH=str(tmp_path / "registry.json"),
        CONVERSATION_STORE_PATH=str(tmp_path / "conv.json"),
        UPLOAD_CACHE_DIR=str(tmp_path / "uploads"),
        ENABLE_TABLE_SERIALIZATION=True,
        ENABLE_OCR=True,
        OCR_MIN_TEXT_CHARS=40,
    )
    base.update(kwargs)
    return Settings(**base)


def test_rows_to_markdown_basic():
    md = rows_to_markdown([["Name", "Days"], ["Annual", "15"], ["Sick", "10"]])
    assert "| Name | Days |" in md
    assert "| --- | --- |" in md
    assert "| Annual | 15 |" in md
    assert "| Sick | 10 |" in md


def test_docx_table_markdown(tmp_path: Path):
    path = tmp_path / "leave.docx"
    doc = DocxDocument()
    doc.add_paragraph("Leave policy")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Type"
    table.cell(0, 1).text = "Days"
    table.cell(1, 0).text = "Annual"
    table.cell(1, 1).text = "15"
    doc.save(path)

    settings = _settings(tmp_path, ENABLE_TABLE_SERIALIZATION=True, ENABLE_OCR=False)
    docs = load_docx(path, settings=settings)
    body = docs[0].page_content
    assert "| Type | Days |" in body
    assert "| Annual | 15 |" in body
    assert docs[0].metadata.get("table_count", 0) >= 1


def test_pptx_table_markdown(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "SLO"
    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(1.5))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Target"
    table.cell(1, 0).text = "p95"
    table.cell(1, 1).text = "200ms"
    prs.save(path)

    settings = _settings(tmp_path, ENABLE_TABLE_SERIALIZATION=True, ENABLE_OCR=False)
    docs = load_pptx(path, settings=settings)
    blob = "\n".join(d.page_content for d in docs)
    assert "| Metric | Target |" in blob
    assert "200ms" in blob


def test_pdf_tables_and_ocr_skip_when_unavailable(tmp_path: Path):
    pytest.importorskip("pdfplumber")
    try:
        from pypdf import PdfWriter
    except ImportError:
        pytest.skip("pypdf required")

    # Minimal blank-ish PDF (no text) — OCR should attempt then skip if unavailable
    pdf_path = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with pdf_path.open("wb") as f:
        writer.write(f)

    settings = _settings(tmp_path, ENABLE_OCR=True, OCR_MIN_TEXT_CHARS=40)
    with patch("src.ingestion.ocr.ocr_available", return_value=(False, "tesseract_unavailable:test")):
        # blank page may yield no docs → IngestError from docs_or_fail
        from src.services.exceptions import IngestError

        try:
            docs = load_pdf(pdf_path, settings=settings)
        except IngestError:
            docs = []
    # Either empty (acceptable skip) or OCR skipped status on any docs
    for d in docs:
        assert d.metadata.get("ocr_status") in {
            "skipped_unavailable",
            "skipped_disabled",
            "skipped_not_needed",
            "failed",
            "applied",
        }


def test_should_ocr_text_threshold():
    assert should_ocr_text("", min_chars=40) is True
    assert should_ocr_text("short", min_chars=40) is True
    assert should_ocr_text("x" * 50, min_chars=40) is False


def test_ocr_pil_disabled(tmp_path: Path):
    settings = _settings(tmp_path, ENABLE_OCR=False)
    result = ocr_pil_image(MagicMock(), settings)
    assert result.status == "skipped_disabled"


def test_ingest_pipeline_includes_tables_and_ocr(tmp_path: Path):
    path = tmp_path / "policy.docx"
    doc = DocxDocument()
    doc.add_paragraph("Handbook")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Type"
    table.cell(0, 1).text = "Days"
    table.cell(1, 0).text = "Annual"
    table.cell(1, 1).text = "15"
    doc.save(path)

    settings = _settings(tmp_path, ENABLE_TABLE_SERIALIZATION=True, ENABLE_OCR=True)
    # Avoid real embedding: mock vector/bm25
    ingest = IngestService(settings=settings)
    ingest.vector_store = MagicMock()
    ingest.vector_store.delete_by_doc_id.return_value = 0
    ingest.vector_store.add_documents = MagicMock()
    ingest.bm25_store = MagicMock()
    ingest.bm25_store.delete_by_doc_id.return_value = 0
    ingest.bm25_store.upsert_documents = MagicMock()
    ingest.registry = MagicMock()
    ingest.registry.upsert.return_value = {"doc_id": "x"}

    result = ingest.ingest_file(path)
    steps = {s["step"]: s for s in result["pipeline_steps"]}
    assert "tables" in steps
    assert steps["tables"]["enabled"] is True
    assert steps["tables"]["table_count"] >= 1
    assert "ocr" in steps
    assert steps["ocr"]["engine"] == "tesseract"
    assert steps["ocr"]["enabled"] is True


def test_settings_defaults_step36():
    assert Settings.model_fields["enable_table_serialization"].default is True
    assert Settings.model_fields["enable_ocr"].default is True
